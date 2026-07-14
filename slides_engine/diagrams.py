"""
slides_engine/diagrams.py
==========================

图表与示意图 —— 让 PPT "会讲数据",而且有设计感。

第一版用 python-pptx 原生图表,是 Office 默认长相,平。这一版全部改成**用原生形状手绘**:
  · 柱/条 —— 圆角渐变柱 + 淡参考线 + 峰值高亮 + 值标签(不是灰扑扑的 Office 柱)
  · 折线 —— 面积渐变填充 + 圆点 marker + 端点值(数据可视化的"设计款")
  · 饼/环 —— 楔形拼色 + 环心大数字高亮 + 右侧色块图例
手绘的好处:设计完全可控,且能和 HTML 预览做到 1:1(彻底摆脱"看不到 pptx 渲染")。
比例永远由几何精确保证 —— 楔形 sweep / 柱高都按真实数值算。
"""
from __future__ import annotations

from math import cos, radians, sin
from typing import List

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ._draw import (
    _grad, _mix, _no_line, _oval, _para, _rect, _rgb, _round, _shadow, _tb,
)
from .styles import DeckStyle


def _palette(style: DeckStyle) -> List[str]:
    a, a2, bg = style.accent, style.accent2, style.bg
    if style.is_dark:
        return [a, a2, _mix(a, "FFFFFF", 0.30), _mix(a2, bg, 0.30),
                _mix(a, bg, 0.48), _mix(a2, "FFFFFF", 0.36)]
    return [a, a2, _mix(a, "000000", 0.20), _mix(a2, "000000", 0.12),
            _mix(a, "FFFFFF", 0.40), _mix(a2, "FFFFFF", 0.42)]


def _grad_round(slide, l, t, w, h, c1, c2, angle=90, radius=0.22):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(l), int(t), int(w), int(h))
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    _grad(sp, c1, c2, angle)
    return sp


def _label(slide, l, t, w, h, text, style, size, color, *, bold=False,
           align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    _, tf = _tb(slide, int(l), int(t), int(w), int(h), anchor)
    _para(tf, text, style, size, color, bold=bold, first=True, align=align, space_after=0,
          line_spacing=1.02)


def _fmt(v) -> str:
    if isinstance(v, float) and v.is_integer():
        v = int(v)
    return str(v)


def draw_chart(slide, style: DeckStyle, s, l, t, w, h):
    if not s.series or not s.series[0].get("values"):
        _label(slide, l, t, w, h, "图表数据缺失 · 用 `标签 | 数值` 或 markdown 表格",
               style, style.pt_body, style.ink_muted)
        return
    ctype = (s.chart_type or "column").lower()
    if ctype in ("pie", "doughnut", "donut"):
        _draw_pie(slide, style, s, l, t, w, h, donut=(ctype != "pie"))
    elif ctype == "bar":
        _draw_bars(slide, style, s, l, t, w, h, horizontal=True)
    elif ctype in ("line", "area"):
        _draw_line(slide, style, s, l, t, w, h, area=(ctype == "area"))
    else:
        _draw_bars(slide, style, s, l, t, w, h, horizontal=False)


# ---------------------------------------------------------------- 柱 / 条 ----
def _draw_bars(slide, style: DeckStyle, s, l, t, w, h, horizontal=False):
    cats = s.categories or [str(i + 1) for i in range(len(s.series[0]["values"]))]
    series = s.series
    multi = len(series) > 1
    pal = _palette(style)
    vals_all = [v for ser in series for v in ser.get("values", [])]
    maxv = max(vals_all) if vals_all else 1
    maxv = maxv or 1

    legend_h = Inches(0.5) if multi else 0
    if multi:
        _legend_row(slide, style, [ser.get("name", f"序列{i+1}") for i, ser in enumerate(series)],
                    pal, l, t, w, Inches(0.36))

    if not horizontal:
        _draw_columns(slide, style, cats, series, pal, maxv, multi,
                      l, t + legend_h, w, h - legend_h)
    else:
        _draw_hbars(slide, style, cats, series, pal, maxv, multi,
                    l, t + legend_h, w, h - legend_h)


def _draw_columns(slide, style, cats, series, pal, maxv, multi, l, t, w, h):
    lab_h = Inches(0.42)          # 底部类别标签
    val_h = Inches(0.34)          # 顶部值标签留白
    plot_t = t + val_h
    plot_h = h - val_h - lab_h
    base_y = plot_t + plot_h
    # 淡参考线(4 条)+ 基线
    for k in range(5):
        gy = plot_t + plot_h * k / 4
        _rect(slide, l, gy, w, Pt(0.75 if k else 1.4),
              style.rule if k else _mix(style.rule, style.ink_muted, 0.4))
    n = len(cats)
    ns = len(series)
    slot = w / n
    group_w = slot * 0.62
    bar_w = group_w / ns
    for ci in range(n):
        slot_l = l + ci * slot + (slot - group_w) / 2
        for si, ser in enumerate(series):
            vals = ser.get("values", [])
            v = vals[ci] if ci < len(vals) else 0
            bh = max(int(plot_h * (v / maxv)), int(Inches(0.03)))
            bx = slot_l + si * bar_w
            c1 = pal[si % len(pal)]
            c2 = _mix(c1, style.bg if style.is_dark else "FFFFFF", 0.30)
            is_peak = (not multi) and abs(v - maxv) < 1e-9
            top_c = c1 if is_peak else _mix(c1, style.bg if style.is_dark else "FFFFFF", 0.12)
            bot_c = _mix(c1, style.bg, 0.35) if style.is_dark else _mix(c1, "000000", 0.06)
            bar = _grad_round(slide, bx + int(bar_w * 0.08), base_y - bh,
                              bar_w * 0.84, bh, top_c, bot_c, angle=90,
                              radius=min(0.5, float(Inches(0.09)) / max(bar_w * 0.84, 1)))
            if is_peak and not style.is_dark:
                _shadow(bar, blur=60000, dist=20000, alpha=22000)
            if not multi:
                _label(slide, bx - int(bar_w * 0.3), base_y - bh - int(val_h), bar_w * 1.6, val_h,
                       _fmt(v), style, style.pt_body, style.accent if is_peak else style.ink_body,
                       bold=True, anchor=MSO_ANCHOR.BOTTOM)
        _label(slide, l + ci * slot, base_y + Inches(0.06), slot, lab_h,
               str(cats[ci]), style, max(style.pt_footnote + 1, 12), style.ink_muted)


def _draw_hbars(slide, style, cats, series, pal, maxv, multi, l, t, w, h):
    ser = series[0]
    vals = ser.get("values", [])
    n = len(cats)
    lab_w = Inches(2.3)
    val_w = Inches(0.9)
    track_l = l + lab_w
    track_w = w - lab_w - val_w
    gap = Inches(0.22)
    row_h = (h - gap * (n - 1)) / n
    bar_h = min(row_h * 0.6, Inches(0.62))
    for i in range(n):
        v = vals[i] if i < len(vals) else 0
        ry = t + i * (row_h + gap) + (row_h - bar_h) / 2
        _label(slide, l, ry - Inches(0.02), lab_w - Inches(0.18), bar_h,
               str(cats[i]), style, style.pt_body, style.ink_body,
               align=PP_ALIGN.RIGHT)
        _round(slide, track_l, ry, track_w, bar_h, style.bg_alt if not style.is_dark
               else _mix(style.bg, "FFFFFF", 0.06), radius=0.5)
        bw = max(int(track_w * (v / maxv)), int(Inches(0.08)))
        is_peak = abs(v - maxv) < 1e-9
        c1 = style.accent if is_peak else _mix(style.accent, style.accent2, 0.5)
        _grad_round(slide, track_l, ry, bw, bar_h, c1, style.accent2, angle=0, radius=0.5)
        _label(slide, track_l + bw + Inches(0.12), ry - Inches(0.02), val_w, bar_h,
               _fmt(v), style, style.pt_body, style.accent if is_peak else style.ink_body,
               bold=True, align=PP_ALIGN.LEFT)


# --------------------------------------------------------------------- 折线 ----
def _draw_line(slide, style: DeckStyle, s, l, t, w, h, area=False):
    cats = s.categories or [str(i + 1) for i in range(len(s.series[0]["values"]))]
    series = s.series
    multi = len(series) > 1
    pal = _palette(style)
    vals_all = [v for ser in series for v in ser.get("values", [])]
    maxv = max(vals_all) if vals_all else 1
    maxv = maxv * 1.12 or 1

    legend_h = Inches(0.5) if multi else 0
    if multi:
        _legend_row(slide, style, [ser.get("name", f"序列{i+1}") for i, ser in enumerate(series)],
                    pal, l, t, w, Inches(0.36))
    lab_h = Inches(0.42)
    val_h = Inches(0.3)
    plot_t = t + legend_h + val_h
    plot_h = h - legend_h - val_h - lab_h
    base_y = plot_t + plot_h
    for k in range(5):
        gy = plot_t + plot_h * k / 4
        _rect(slide, l, gy, w, Pt(0.75 if k else 1.4),
              style.rule if k else _mix(style.rule, style.ink_muted, 0.4))
    n = len(cats)
    step = w / max(n - 1, 1) if n > 1 else w
    xs = [l + (i * step if n > 1 else w / 2) for i in range(n)]

    for si, ser in enumerate(series):
        vals = ser.get("values", [])
        pts = [(int(xs[i]), int(base_y - plot_h * (min(vals[i], maxv) / maxv)))
               for i in range(min(n, len(vals)))]
        if len(pts) < 1:
            continue
        col = pal[si % len(pal)]
        if area and not multi and len(pts) >= 2:
            try:
                fb = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=1.0)
                rest = pts[1:] + [(pts[-1][0], int(base_y)), (pts[0][0], int(base_y))]
                fb.add_line_segments(rest, close=True)
                shp = fb.convert_to_shape()
                _grad(shp, _mix(col, style.bg, 0.55), style.bg, angle=90)
            except Exception:
                pass
        for a, b in zip(pts, pts[1:]):
            cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, a[0], a[1], b[0], b[1])
            cn.line.color.rgb = _rgb(col); cn.line.width = Pt(3)
        mr = int(Inches(0.11))
        for i, (px, py) in enumerate(pts):
            _oval(slide, px - mr, py - mr, mr * 2, mr * 2, col)
            _oval(slide, px - mr // 2, py - mr // 2, mr, mr, style.bg)
            if not multi:
                _label(slide, px - int(Inches(0.6)), py - int(val_h) - int(Inches(0.12)),
                       int(Inches(1.2)), val_h, _fmt(vals[i]), style,
                       max(style.pt_footnote + 1, 12), style.ink_body, bold=True,
                       anchor=MSO_ANCHOR.BOTTOM)
    for i in range(n):
        _label(slide, xs[i] - step / 2, base_y + Inches(0.06), step, lab_h,
               str(cats[i]), style, max(style.pt_footnote + 1, 12), style.ink_muted)


# ---------------------------------------------------------------- 饼 / 环 ----
def _draw_pie(slide, style: DeckStyle, s, l, t, w, h, donut=False):
    labels = s.categories or [str(i + 1) for i in range(len(s.series[0]["values"]))]
    vals = list(s.series[0].get("values", []))
    total = sum(vals) or 1
    pal = _palette(style)
    side = int(min(h, w * 0.52))
    px = int(l)
    py = int(t + (h - side) / 2)

    # MSO_SHAPE.PIE 的 adj1/adj2 是 ST_AdjAngle(60000 分之一度)· 但 python-pptx 的
    # adjustments setter 一律按 value*100000 落 XML。 所以要把"度"乘 60000/100000=0.6 再塞,
    # 否则每个角被放大 1.667 倍 → 弧全错(饼图对不上数据的真凶)。
    _A = 60000.0 / 100000.0
    start = 270.0                       # 12 点起 · 顺时针
    nz = [v for v in vals if v > 0]
    peak_i = max(range(len(vals)), key=lambda i: vals[i]) if vals else 0
    for i, v in enumerate(vals):
        if v <= 0:
            continue
        sweep = 360.0 * (v / total)
        end = start + sweep
        if len(nz) == 1 or sweep >= 359.9:          # 单段 100% → 整圆(避免退化 wedge)
            wedge = _oval(slide, px, py, side, side, pal[i % len(pal)])
        else:
            wedge = slide.shapes.add_shape(MSO_SHAPE.PIE, px, py, side, side)
            try:
                wedge.adjustments[0] = (start % 360) * _A
                wedge.adjustments[1] = (end % 360) * _A
            except Exception:
                pass
            wedge.fill.solid(); wedge.fill.fore_color.rgb = _rgb(pal[i % len(pal)])
        wedge.line.color.rgb = _rgb(style.bg); wedge.line.width = Pt(2.25)
        wedge.shadow.inherit = False
        start = end
    if donut:
        hole = int(side * 0.56)
        hx = px + (side - hole) // 2
        hy = py + (side - hole) // 2
        _oval(slide, hx, hy, hole, hole, style.bg)
        pct = round(100 * vals[peak_i] / total) if vals else 0
        _label(slide, hx, hy + int(hole * 0.16), hole, int(hole * 0.42),
               f"{pct}%", style, style.pt_kpi, style.accent, bold=True)
        _label(slide, hx, hy + int(hole * 0.56), hole, int(hole * 0.26),
               str(labels[peak_i]) if labels else "", style, style.pt_footnote,
               style.ink_muted)
    # 右侧图例
    lx = px + side + int(Inches(0.7))
    lw = int(l + w - lx)
    rows = len(labels)
    rh = min(int(Inches(0.62)), int((h) / max(rows, 1)))
    ly0 = int(t + (h - rh * rows) / 2)
    chip = int(Inches(0.26))
    for i, lab in enumerate(labels):
        ry = ly0 + i * rh
        _round(slide, lx, ry + (rh - chip) // 2, chip, chip, pal[i % len(pal)], radius=0.35)
        pct = round(100 * vals[i] / total) if i < len(vals) else 0
        _label(slide, lx + chip + int(Inches(0.16)), ry, lw - chip - int(Inches(0.9)), rh,
               str(lab), style, style.pt_body, style.ink_body, align=PP_ALIGN.LEFT)
        _label(slide, lx + lw - int(Inches(0.85)), ry, int(Inches(0.85)), rh,
               f"{pct}%", style, style.pt_body, style.accent, bold=True, align=PP_ALIGN.RIGHT)


def _legend_row(slide, style: DeckStyle, names, pal, l, t, w, h):
    chip = int(Inches(0.2))
    x = int(l)
    y = int(t)
    for i, nm in enumerate(names):
        _round(slide, x, y + (int(h) - chip) // 2, chip, chip, pal[i % len(pal)], radius=0.35)
        tw = int(Inches(0.16)) + max(int(Inches(0.9)), len(str(nm)) * int(Inches(0.14)))
        _label(slide, x + chip + int(Inches(0.1)), y, tw, int(h), str(nm), style,
               style.pt_footnote + 1, style.ink_muted, align=PP_ALIGN.LEFT)
        x += chip + int(Inches(0.1)) + tw + int(Inches(0.25))


# --------------------------------------------------------------------- 流程 ----
def draw_flow(slide, style: DeckStyle, steps: List[dict], l, t, w, h):
    """横向流程图:圆角步骤框 + 序号圆点 + 箭头。 最多 5 步(超出收字)。"""
    labels = [b.get("text", "") for b in steps][:5]
    n = len(labels)
    if n == 0:
        return
    gap = Inches(0.5)
    bw = (w - gap * (n - 1)) / n
    bh = min(Inches(2.1), h)
    bt = t + (h - bh) / 2
    fs = style.pt_body if n <= 3 else max(style.pt_body - 3, 12)
    for i, label in enumerate(labels):
        bl = l + i * (bw + gap)
        card = _round(slide, bl, bt, bw, bh, style.bg_alt,
                      line=(style.rule if style.is_dark else None))
        if not style.is_dark:
            _shadow(card)
        _rect(slide, bl, bt, bw, Pt(6), style.accent)
        chip = int(Inches(0.5))
        cx = int(bl) + (int(bw) - chip) // 2
        _oval(slide, cx, int(bt) + int(Inches(0.34)), chip, chip, style.accent)
        _, tfn = _tb(slide, cx, int(bt) + int(Inches(0.34)), chip, chip, MSO_ANCHOR.MIDDLE)
        _para(tfn, str(i + 1), style, style.pt_heading, style.on_accent, bold=True, first=True,
              align=PP_ALIGN.CENTER, space_after=0)
        _, tf = _tb(slide, int(bl) + int(Inches(0.22)), int(bt) + int(Inches(1.05)),
                    int(bw) - int(Inches(0.44)), int(bh) - int(Inches(1.15)), MSO_ANCHOR.TOP)
        _para(tf, label, style, fs, style.ink_body, first=True, align=PP_ALIGN.CENTER,
              space_after=0, line_spacing=1.12)
        if i < n - 1:
            ar_w = int(gap)
            ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, int(bl) + int(bw) + int(Inches(0.06)),
                                        int(bt) + int(bh) // 2 - int(Inches(0.16)),
                                        ar_w - int(Inches(0.12)), int(Inches(0.32)))
            ar.fill.solid(); ar.fill.fore_color.rgb = _rgb(style.accent)
            _no_line(ar)
