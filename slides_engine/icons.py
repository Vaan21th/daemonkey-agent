"""
slides_engine/icons.py
=======================

单色线性图标 —— 商务 PPT 的设计感,一半在图标。

pptx 塞不进 SVG 线性图标(python-pptx 无贝塞尔曲线 API),所以这里用**原生形状 + 可靠内置图形**
(齿轮 GEAR_6 / 五角星 / 箭头 / 圆环 / 细条)拼出一套克制的单色图标。好处:
  · 完全矢量、任意缩放不糊
  · 一键回色(套主题紫)
  · 和 HTML 预览用同一套几何 → 所见即所得

每个图标画在正方框 (x, y, sz) 里,统一留白、线宽随尺寸缩放。未知名字兜底成"圆点"。
"""
from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from ._draw import _oval, _para, _rect, _rgb, _round, _tb

_EMU_PT = 12700.0

# 语义别名 → 规范名
_ALIAS = {
    "goal": "target", "aim": "target", "bullseye": "target",
    "data": "chart", "analytics": "chart", "kpi": "chart", "metric": "chart",
    "up": "growth", "trend": "growth", "rise": "growth", "increase": "growth",
    "bulb": "idea", "insight": "idea", "innovation": "idea", "creative": "idea",
    "process": "gear", "settings": "gear", "config": "gear", "system": "gear", "engine": "gear",
    "team": "people", "user": "people", "users": "people", "audience": "people", "client": "people",
    "revenue": "money", "cash": "money", "profit": "money", "cost": "money", "price": "money",
    "launch": "rocket", "ship": "rocket", "fast": "rocket", "start": "rocket",
    "report": "doc", "file": "doc", "document": "doc", "content": "doc", "article": "doc",
    "done": "check", "ok": "check", "verify": "check", "quality": "check",
    "highlight": "star", "feature": "star", "premium": "star", "best": "star",
    "clock": "time", "schedule": "time", "timeline": "time", "efficiency": "time",
    "find": "search", "research": "search", "discover": "search", "radar": "search",
    "milestone": "flag", "goalpost": "flag",
    "connect": "link", "integration": "link", "chain": "link",
}


class _GlyphStyle:
    font_en = "Segoe UI"
    font_cjk = "Microsoft YaHei"


_GS = _GlyphStyle()


def resolve(name: str) -> str:
    n = (name or "").strip().lower()
    return _ALIAS.get(n, n)


def has(name: str) -> bool:
    return resolve(name) in _ICONS


def _pt(emu: float) -> float:
    return max(emu / _EMU_PT, 1.4)


def _ring(slide, x, y, d, ink, wpt):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(x), int(y), int(d), int(d))
    sp.fill.background()
    sp.line.color.rgb = _rgb(ink); sp.line.width = Pt(wpt)
    sp.shadow.inherit = False
    return sp


def _auto(slide, shp, x, y, w, h, ink, *, fill=True, wpt=2.4):
    sp = slide.shapes.add_shape(shp, int(x), int(y), int(w), int(h))
    if fill:
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(ink); sp.line.fill.background()
    else:
        sp.fill.background(); sp.line.color.rgb = _rgb(ink); sp.line.width = Pt(wpt)
    sp.shadow.inherit = False
    return sp


def _glyph(slide, x, y, w, h, ch, ink, pt):
    _, tf = _tb(slide, int(x), int(y), int(w), int(h), MSO_ANCHOR.MIDDLE)
    _para(tf, ch, _GS, pt, ink, bold=True, first=True, align=PP_ALIGN.CENTER, space_after=0)


def draw_icon(slide, name, x, y, sz, ink, style=None):
    """在正方框 (x, y, sz) 里画一个单色图标(ink=颜色)。"""
    x, y, sz = int(x), int(y), int(sz)
    n = resolve(name)
    lw = max(sz * 0.06, 1)
    m = sz * 0.14
    ix, iy, iw = x + m, y + m, sz - 2 * m
    fn = _ICONS.get(n)
    if fn:
        fn(slide, x, y, sz, ix, iy, iw, lw, ink)
    else:
        d = sz * 0.34
        _oval(slide, x + (sz - d) / 2, y + (sz - d) / 2, d, d, ink)


def _icon_target(s, x, y, sz, ix, iy, iw, lw, ink):
    for k in (1.0, 0.6):
        d = iw * k
        _ring(s, ix + (iw - d) / 2, iy + (iw - d) / 2, d, ink, _pt(lw) * 1.4)
    dd = iw * 0.2
    _oval(s, ix + (iw - dd) / 2, iy + (iw - dd) / 2, dd, dd, ink)


def _icon_chart(s, x, y, sz, ix, iy, iw, lw, ink):
    gap = iw * 0.14
    bw = (iw - gap * 2) / 3
    base = iy + iw
    for i, hf in enumerate((0.45, 0.72, 1.0)):
        bh = iw * hf
        _round(s, ix + i * (bw + gap), base - bh, bw, bh, ink, radius=0.28)


def _icon_growth(s, x, y, sz, ix, iy, iw, lw, ink):
    _auto(s, MSO_SHAPE.BENT_UP_ARROW, ix, iy + iw * 0.12, iw, iw * 0.8, ink, fill=True)


def _icon_idea(s, x, y, sz, ix, iy, iw, lw, ink):
    bd = iw * 0.62
    _ring(s, ix + (iw - bd) / 2, iy, bd, ink, _pt(lw) * 1.5)
    bw = iw * 0.24
    _round(s, ix + (iw - bw) / 2, iy + bd * 0.84, bw, iw * 0.2, ink, radius=0.4)
    _rect(s, ix + (iw - bw) / 2, iy + bd * 1.04, bw, max(iw * 0.05, 1), ink)


def _icon_gear(s, x, y, sz, ix, iy, iw, lw, ink):
    _auto(s, MSO_SHAPE.GEAR_6, ix, iy, iw, iw, ink, fill=True)
    dd = iw * 0.34
    _oval(s, ix + (iw - dd) / 2, iy + (iw - dd) / 2, dd, dd, "FFFFFF")


def _icon_people(s, x, y, sz, ix, iy, iw, lw, ink):
    pw = iw * 0.52

    def person(px):
        hd = pw * 0.42
        _oval(s, px + (pw - hd) / 2, iy + iw * 0.06, hd, hd, ink)
        _auto(s, MSO_SHAPE.CHORD, px, iy + iw * 0.44, pw, pw * 0.9, ink, fill=True)
    person(ix)
    person(ix + iw - pw)


def _icon_money(s, x, y, sz, ix, iy, iw, lw, ink):
    _ring(s, ix, iy, iw, ink, _pt(lw) * 1.5)
    _glyph(s, ix, iy, iw, iw, "\u00a5", ink, max(int(iw / _EMU_PT * 0.82), 12))


def _icon_rocket(s, x, y, sz, ix, iy, iw, lw, ink):
    bw = iw * 0.44
    bx = ix + (iw - bw) / 2
    _auto(s, MSO_SHAPE.ISOSCELES_TRIANGLE, bx, iy, bw, iw * 0.34, ink, fill=True)
    _round(s, bx, iy + iw * 0.26, bw, iw * 0.48, ink, radius=0.4)
    _auto(s, MSO_SHAPE.ISOSCELES_TRIANGLE, ix, iy + iw * 0.6, bw * 0.5, iw * 0.32, ink, fill=True)
    _auto(s, MSO_SHAPE.ISOSCELES_TRIANGLE, ix + iw - bw * 0.5, iy + iw * 0.6,
          bw * 0.5, iw * 0.32, ink, fill=True)


def _icon_doc(s, x, y, sz, ix, iy, iw, lw, ink):
    pw, ph = iw * 0.72, iw
    px = ix + (iw - pw) / 2
    _round(s, px, iy, pw, ph, ink, radius=0.08)
    for i in range(3):
        _rect(s, px + pw * 0.18, iy + ph * (0.3 + i * 0.19), pw * 0.64, max(iw * 0.05, 1), "FFFFFF")


def _icon_check(s, x, y, sz, ix, iy, iw, lw, ink):
    _oval(s, ix, iy, iw, iw, ink)
    _glyph(s, ix, iy, iw, iw, "\u2713", "FFFFFF", max(int(iw / _EMU_PT * 0.72), 12))


def _icon_star(s, x, y, sz, ix, iy, iw, lw, ink):
    _auto(s, MSO_SHAPE.STAR_5_POINT, ix, iy, iw, iw, ink, fill=True)


def _icon_time(s, x, y, sz, ix, iy, iw, lw, ink):
    _ring(s, ix, iy, iw, ink, _pt(lw) * 1.5)
    cx, cy = ix + iw / 2, iy + iw / 2
    _rect(s, cx - max(iw * 0.03, 1), cy - iw * 0.3, max(iw * 0.06, 1), iw * 0.3, ink)
    _rect(s, cx, cy - max(iw * 0.03, 1), iw * 0.24, max(iw * 0.06, 1), ink)


def _icon_search(s, x, y, sz, ix, iy, iw, lw, ink):
    d = iw * 0.72
    _ring(s, ix, iy, d, ink, _pt(lw) * 1.6)
    h = _round(s, ix + d * 0.76, iy + d * 0.76, iw * 0.34, max(iw * 0.11, 2), ink, radius=0.5)
    h.rotation = 45


def _icon_flag(s, x, y, sz, ix, iy, iw, lw, ink):
    _rect(s, ix + iw * 0.14, iy, max(iw * 0.09, 2), iw, ink)
    _auto(s, MSO_SHAPE.PENTAGON, ix + iw * 0.2, iy + iw * 0.04, iw * 0.62, iw * 0.42, ink, fill=True)


def _icon_link(s, x, y, sz, ix, iy, iw, lw, ink):
    rw, rh = iw * 0.56, iw * 0.4
    _round(s, ix, iy + iw * 0.1, rw, rh, ink, radius=0.5)
    _round(s, ix + iw - rw, iy + iw - rh - iw * 0.1, rw, rh, ink, radius=0.5)
    _rect(s, ix + iw * 0.42, iy + iw * 0.42, iw * 0.16, iw * 0.16, ink)


_ICONS = {
    "target": _icon_target, "chart": _icon_chart, "growth": _icon_growth,
    "idea": _icon_idea, "gear": _icon_gear, "people": _icon_people,
    "money": _icon_money, "rocket": _icon_rocket, "doc": _icon_doc,
    "check": _icon_check, "star": _icon_star, "time": _icon_time,
    "search": _icon_search, "flag": _icon_flag, "link": _icon_link,
}
