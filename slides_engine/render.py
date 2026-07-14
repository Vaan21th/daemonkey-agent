"""
slides_engine/render.py
========================

编排层:Slide 列表 → .pptx。 绘图原语在 _draw.py · 版式在 layouts.py。
高级感来源见那两个文件;这里只负责建 deck、按 layout 分发、落盘。
"""
from __future__ import annotations

from pathlib import Path
from typing import List, Optional

from pptx import Presentation

from . import layouts as L
from . import fonts as _fonts
from ._draw import SLIDE_H, SLIDE_W, set_active_fonts
from .parse import Slide
from .styles import DeckStyle, get_style

# 这些版式自带整页构图 · 不叠页脚页码
_NO_FOOTER = {"cover", "section", "closing"}


def _cover_slide(prs, style: DeckStyle, cover: dict, here=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    s = Slide(layout="cover", title=cover.get("title"), subtitle=cover.get("subtitle"),
              kicker=cover.get("audience"), footer=cover.get("footer"),
              image=cover.get("image"), image_prompt=cover.get("prompt"))
    L.cover(slide, s, style, here, cover.get("cover_layout", "auto"))


def render_deck(slides: List[Slide], output_path, *, cover: Optional[dict] = None,
                style: str | DeckStyle = "light_studio", here_dir=None) -> Path:
    st = style if isinstance(style, DeckStyle) else get_style(style)
    here = Path(here_dir) if here_dir else None
    # 读系统字体库 · 有设计字体就自动用上 · 再按 font_role 偏选(mono/hand/serif)
    set_active_fonts(_fonts.pick_fonts_for(getattr(st, "font_role", "sans")))

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    if cover and cover.get("title"):
        _cover_slide(prs, st, cover, here)

    page = 0
    sec_no = 0        # 章节序号 · 给 section 水印当巨型数字(替代丑陋的 §)
    img_no = 0        # 配图页序号 · 给 image 版式左右交替(打破"图永远在右")
    for s in slides:
        slide = prs.slides.add_slide(blank)
        lay = s.layout
        if lay in _NO_FOOTER or lay == "statement":
            page_num = None
        else:
            page += 1
            page_num = page

        if lay == "cover":
            L.cover(slide, s, st, here)
        elif lay == "section":
            sec_no += 1
            L.section(slide, s, st, sec_no)
        elif lay == "image":
            img_no += 1
            L.image(slide, s, st, page_num, here, img_no)
        elif lay == "statement":
            L.statement(slide, s, st, page_num)
        elif lay == "two_col":
            L.two_col(slide, s, st, page_num)
        elif lay == "metrics":
            L.metrics(slide, s, st, page_num)
        elif lay == "pillars":
            L.pillars(slide, s, st, page_num)
        elif lay == "chart":
            L.chart(slide, s, st, page_num)
        elif lay == "flow":
            L.flow(slide, s, st, page_num)
        elif lay == "sources":
            L.sources(slide, s, st, page_num)
        elif lay == "closing":
            L.closing(slide, s, st)
        else:
            L.bullets(slide, s, st, page_num)

        if s.notes:
            slide.notes_slide.notes_text_frame.text = s.notes

    return _save(prs, Path(output_path))


def _save(prs, output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    target = output_path
    for i in range(2, 6):
        try:
            prs.save(str(target))
            return target
        except PermissionError:
            target = output_path.with_stem(f"{output_path.stem}-v{i}")
    prs.save(str(target))
    return target
