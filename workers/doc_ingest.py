"""
workers/doc_ingest.py
=====================

把本地文档抽成统一 markdown 文本 —— 知识库入库的第一步(眼睛)。

支持: .md/.markdown/.txt 直读 · .docx(python-docx) · .pptx(python-pptx) · .pdf(pypdf 文本型)
不做: OCR(扫描件返回空并提示) · 图片 · 复杂表格结构化(表格按行拼可读文本)

设计取向:尽量保留 / 造出 `## 小节` 标题(PDF 按页 · PPTX 按幻灯片),
让 memory_index._chunk_markdown 能按节切块 · 检索命中时能 cite 回"第 N 页 / 第 N 页幻灯片"。

复用现有地基:PDF 走 pypdf(与 agent_tools/pdf_read.py 同款)· 不重造轮子。
"""

from __future__ import annotations

import re
from pathlib import Path

SUPPORTED_EXT = {".md", ".markdown", ".txt", ".text", ".docx", ".pptx", ".pdf"}
MAX_BYTES = 50 * 1024 * 1024  # 50MB 上限,和 pdf_read 一致


class IngestError(Exception):
    """抽取失败(格式不支持 / 缺依赖 / 扫描件无文字 / 解析炸)——调用方给可读提示。"""


def _clean(raw: str) -> str:
    """折叠多余空白 · 去行首缩进 · 三连空行压成两个。"""
    raw = re.sub(r"[ \t]+", " ", raw)
    raw = re.sub(r"\n[ \t]+", "\n", raw)
    raw = re.sub(r"\n{3,}", "\n\n", raw)
    return raw.strip()


def _from_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _from_docx(path: Path) -> str:
    import docx  # python-docx

    doc = docx.Document(str(path))
    lines: list[str] = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style = (para.style.name if para.style else "") or ""
        if style.lower().startswith("heading"):
            digits = "".join(ch for ch in style if ch.isdigit())
            level = min(int(digits), 4) if digits else 2
            lines.append(f"\n{'#' * level} {text}")
        else:
            lines.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            if any(cells):
                lines.append(" | ".join(cells))
    return "\n".join(lines)


def _from_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"\n## 第 {i} 页幻灯片")
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    out.append(text)
    return "\n".join(out)


def _from_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception as exc:  # noqa: BLE001 — 任何解密失败都当加密件挡回
            raise IngestError("PDF 加密/有密码 · 知识库不支持自动破解") from exc
    out: list[str] = []
    for i, page in enumerate(reader.pages, 1):
        try:
            text = page.extract_text() or ""
        except Exception:  # noqa: BLE001 — 单页炸不连累整篇
            text = ""
        text = _clean(text)
        out.append(f"\n## 第 {i} 页\n{text}")
    return "\n".join(out)


def _dispatch(ext: str, path: Path) -> tuple[str, str]:
    """按扩展名分发抽取 · 返回 (text, doc_type)。"""
    if ext in (".md", ".markdown"):
        return _from_text(path), "md"
    if ext in (".txt", ".text"):
        return _from_text(path), "txt"
    if ext == ".docx":
        return _from_docx(path), "docx"
    if ext == ".pptx":
        return _from_pptx(path), "pptx"
    if ext == ".pdf":
        return _from_pdf(path), "pdf"
    raise IngestError(f"暂不支持的格式 {ext}")


def extract(path: str | Path) -> tuple[str, str, str]:
    """抽本地文档为 markdown 文本。

    Returns: (title, markdown_text, doc_type)
    Raises: IngestError —— 文件不存在 / 格式不支持 / 过大 / 缺依赖 / 无文字(扫描件)。
    """
    p = Path(path)
    if not p.is_absolute():
        p = Path.cwd() / p
    if not p.exists() or not p.is_file():
        raise IngestError(f"文件不存在或不是文件: {p}")

    ext = p.suffix.lower()
    if ext not in SUPPORTED_EXT:
        raise IngestError(
            f"暂不支持 {ext} · 目前支持 md / txt / docx / pptx / pdf(文本型)"
        )
    if p.stat().st_size > MAX_BYTES:
        raise IngestError(f"文件过大 (>{MAX_BYTES // 1024 // 1024}MB) · 拆开后再灌")

    try:
        text, doc_type = _dispatch(ext, p)
    except IngestError:
        raise
    except ImportError as exc:
        raise IngestError(f"缺少解析依赖: {exc} · pip install 后重试") from exc
    except Exception as exc:  # noqa: BLE001 — 收敛成可读错误交给上层
        raise IngestError(f"解析失败: {type(exc).__name__}: {exc}") from exc

    text = _clean(text)
    if len(text.strip()) < 5:
        raise IngestError(
            "没抽到文字 · 大概率是扫描件 / 图片型文档 · 需要 OCR(暂未支持)"
        )
    return p.stem, text, doc_type
