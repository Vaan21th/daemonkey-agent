"""
slides_engine/_draw.py
=======================

升级版绘图原语 —— "设计感"的地基。

第一版只有白底 + 文本框 + 细线,所以平。这一版补齐真正拉开档次的东西:
  · 渐变色块 / 满版渐变底(depth)
  · 圆角卡片 + 柔和投影(卡片浮起来)
  · 彩色方块 marker(不是灰扑扑的 em-dash)
  · 章节巨型数字水印(编辑设计里最"高级"的一招)
  · 统一 header(眉标 hairline 在标题之上 · 彻底根治眉标压标题的 bug)
  · 装饰圆(非对称构图的点睛)
纯几何 + python-pptx · 供 layouts.py 拼装。
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import List, Optional, Tuple

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from .styles import DeckStyle

_EMU_PER_IN = 914400.0

# inline markdown:LLM 常把正文当普通 markdown 写(**粗体**、`码`、[文字](链接)),
# 不解析就整串字面量漏进 PPT(卷七十九事故:**合作模式** 出现在幻灯片正文)。
_LINK_RE = re.compile(r"\[([^\]]+)\]\((?:[^)]*)\)")
_CODE_RE = re.compile(r"`([^`]+)`")
_BOLD_RE = re.compile(r"\*\*(.+?)\*\*|__(.+?)__")
_STRAY_HEAD_RE = re.compile(r"^\s{0,3}#{1,6}\s+")     # 漏进来的行首 ##
_ITALIC_RE = re.compile(r"(?<![\*_])[\*_](?![\*_\s])(.+?)(?<![\*_\s])[\*_](?![\*_])")
_ICON_TOKEN_RE = re.compile(r"\s*\{[a-z][a-z_]{2,18}\}\s*")   # 残留图标 token(如正文里漏的 {search})


def _clean_inline(text: str) -> str:
    t = _STRAY_HEAD_RE.sub("", text or "")
    t = _LINK_RE.sub(r"\1", t)      # 链接只留可读文字
    t = _CODE_RE.sub(r"\1", t)      # 反引号去壳
    t = _ITALIC_RE.sub(r"\1", t)    # 单星号斜体去壳(不加斜体·避免误伤)
    t = _ICON_TOKEN_RE.sub(" ", t).strip()   # 清残留 {icon} 字面量·图标只该被版式消费·不该出现在文字里
    return t


def _inline_segments(text: str) -> List[Tuple[str, bool]]:
    """把一段可能含 **粗体** 的文本切成 [(片段, 是否加粗)] · 供多 run 渲染。"""
    t = _clean_inline(text)
    segs: List[Tuple[str, bool]] = []
    pos = 0
    for m in _BOLD_RE.finditer(t):
        if m.start() > pos:
            segs.append((t[pos:m.start()], False))
        segs.append((m.group(1) or m.group(2) or "", True))
        pos = m.end()
    if pos < len(t):
        segs.append((t[pos:], False))
    return segs or [("", False)]


def _text_units(s: str) -> float:
    """估算文本宽度(以'一个全角字 = 1 单位'计)· 给 _fit_pt 收缩字号用。"""
    u = 0.0
    for ch in (s or ""):
        o = ord(ch)
        if o > 0x2E7F or o in (0x2014, 0x2013):     # CJK/全角/破折号
            u += 1.0
        elif ch in " iIl.,:;'!|()":
            u += 0.34
        else:
            u += 0.56
    return u


def _fit_pt(text: str, base_pt, box_w_emu, *, max_lines: int = 1,
            min_pt: int = 14, safety: float = 0.94) -> int:
    """按可用宽度把字号从 base 往下收,保证 max_lines 行内塞得下(防大字溢出/重叠)。"""
    u = _text_units(text)
    if u <= 0:
        return int(base_pt)
    cap_in = (float(box_w_emu) / _EMU_PER_IN) * max_lines * safety
    fit = cap_in * 72.0 / u
    return int(max(min_pt, min(float(base_pt), fit)))

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.75)
CONTENT_W = SLIDE_W - MARGIN * 2
BODY_TOP = Inches(2.25)   # 内页正文统一起点(重复=对齐=高级)


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6.replace("#", "").upper())


def _mix(a: str, b: str, t: float) -> str:
    """按 t∈[0,1] 混合两个 hex 颜色(t=0→a · t=1→b)。做渐变/淡色调用。"""
    a = a.replace("#", ""); b = b.replace("#", "")
    ar, ag, ab = int(a[0:2], 16), int(a[2:4], 16), int(a[4:6], 16)
    br, bg, bb = int(b[0:2], 16), int(b[2:4], 16), int(b[4:6], 16)
    r = round(ar + (br - ar) * t); g = round(ag + (bg - ag) * t); bl = round(ab + (bb - ab) * t)
    return f"{r:02X}{g:02X}{bl:02X}"


def _bg(slide, style: DeckStyle, color: Optional[str] = None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(color or style.bg)
    _texture(slide, style)


def _no_line(sp):
    sp.line.fill.background()
    sp.shadow.inherit = False


def _grad(sp, c1: str, c2: str, angle: int = 45):
    sp.fill.gradient()
    stops = sp.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = _rgb(c1)
    stops[1].position = 1.0
    stops[1].color.rgb = _rgb(c2)
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    _no_line(sp)


def _rect(slide, l, t, w, h, fill, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(l), int(t), int(w), int(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    if line:
        sp.line.color.rgb = _rgb(line); sp.line.width = line_w or Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _grad_rect(slide, l, t, w, h, c1, c2, angle=45):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(l), int(t), int(w), int(h))
    _grad(sp, c1, c2, angle)
    return sp


def _set_gs_alpha(sp, alphas):
    """给渐变各 stop 注入透明度(%)· 做"暗→透"遮罩用。"""
    spPr = sp._element.spPr
    gf = spPr.find(qn("a:gradFill"))
    if gf is None:
        return
    gsl = gf.find(qn("a:gsLst"))
    if gsl is None:
        return
    stops = gsl.findall(qn("a:gs"))
    for gs, a in zip(stops, alphas):
        clr = gs.find(qn("a:srgbClr"))
        if clr is None:
            continue
        for old in clr.findall(qn("a:alpha")):
            clr.remove(old)
        v = int(max(0, min(100, a)) * 1000)
        clr.append(clr.makeelement(qn("a:alpha"), {"val": str(v)}))


def _scrim(slide, l, t, w, h, color, a1=85, a2=0, angle=0):
    """渐变遮罩(a1→a2 透明度)· 压在满版图上保证文字可读 · angle=0 左暗右透。"""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(l), int(t), int(w), int(h))
    sp.fill.gradient()
    gs = sp.fill.gradient_stops
    gs[0].position = 0.0; gs[0].color.rgb = _rgb(color)
    gs[1].position = 1.0; gs[1].color.rgb = _rgb(color)
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    _set_gs_alpha(sp, [a1, a2])
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def _round(slide, l, t, w, h, fill, radius=0.055, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(l), int(t), int(w), int(h))
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    if line:
        sp.line.color.rgb = _rgb(line); sp.line.width = line_w or Pt(1)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _oval(slide, l, t, w, h, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, int(l), int(t), int(w), int(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    _no_line(sp)
    return sp


def _corner_marks(slide, style: DeckStyle, color: Optional[str] = None):
    """四角 L 形取景框标记(科技/赛博感)· decor='corner' 用。"""
    c = color or style.accent
    arm = int(Inches(0.72)); th = int(Pt(3.2)); m = int(Inches(0.4))
    W, H = int(SLIDE_W), int(SLIDE_H)
    _rect(slide, m, m, arm, th, c); _rect(slide, m, m, th, arm, c)                          # 左上
    _rect(slide, W - m - arm, m, arm, th, c); _rect(slide, W - m - th, m, th, arm, c)        # 右上
    _rect(slide, m, H - m - th, arm, th, c); _rect(slide, m, H - m - arm, th, arm, c)        # 左下
    _rect(slide, W - m - arm, H - m - th, arm, th, c); _rect(slide, W - m - th, H - m - arm, th, arm, c)  # 右下


def _shadow(sp, *, blur=95000, dist=32000, direction=5400000, alpha=17000, color="222634"):
    """给形状加柔和外投影(卡片浮起感)· 标准 DrawingML outerShdw。"""
    spPr = sp._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shd = eff.makeelement(qn("a:outerShdw"), {
        "blurRad": str(blur), "dist": str(dist), "dir": str(direction), "rotWithShape": "0"})
    clr = shd.makeelement(qn("a:srgbClr"), {"val": color.replace("#", "").upper()})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(alpha)}))
    shd.append(clr); eff.append(shd); spPr.append(eff)


# ─────────────────────────────────────────────────────────────
# 结构级设计 token 的绘图原语 —— 半透明 / 发光 / 纹理 / 统一卡片
# 让"同一渲染器 + 不同 token = 不同艺术方向"成立(玻璃/故障/手绘)
# ─────────────────────────────────────────────────────────────
def _apply_alpha(sp, alpha_pct: int):
    """给纯色填充加透明度(0-100)· 玻璃拟态的磨砂面板靠它。"""
    spPr = sp._element.spPr
    sf = spPr.find(qn("a:solidFill"))
    if sf is None:
        return
    clr = sf.find(qn("a:srgbClr"))
    if clr is None:
        return
    for a in clr.findall(qn("a:alpha")):
        clr.remove(a)
    v = int(max(0, min(100, alpha_pct)) * 1000)   # a:alpha 单位 = 千分之一百分比
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(v)}))


def _glow(sp, color: str, rad: int = 165000, alpha: int = 44000):
    """霓虹发光(真 a:glow)· 故障/暗色发布会的边缘辉光。"""
    spPr = sp._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    g = eff.makeelement(qn("a:glow"), {"rad": str(rad)})
    clr = g.makeelement(qn("a:srgbClr"), {"val": color.replace("#", "").upper()})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": str(alpha)}))
    g.append(clr); eff.append(g); spPr.append(eff)


def _panel_border(style: DeckStyle):
    """卡片描边(按 stroke_style)· 返回 (hex, width) 或 None。 clean=现状(深底描边/浅底无边)。"""
    ss = getattr(style, "stroke_style", "clean")
    if ss == "hairline":
        return style.rule, Pt(0.75)
    if ss == "hard":
        return style.accent, Pt(2.0)
    if ss == "sketch":
        return style.ink_muted, Pt(1.6)
    return (style.rule, Pt(1.0)) if style.is_dark else None


def _panel(slide, style: DeckStyle, l, t, w, h, *, fill: Optional[str] = None, radius=None):
    """统一卡片:圆角 + 透明度 + 渐变 + 描边 + 阴影/发光 —— 全读 style token。
    默认 token 下等价于旧「_round(bg_alt) + 浅底 _shadow」(零回归)。"""
    fill = fill or style.bg_alt
    rad = style.corner_radius if radius is None else radius
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(l), int(t), int(w), int(h))
    try:
        sp.adjustments[0] = rad
    except Exception:
        pass
    if getattr(style, "panel_gradient", False):
        c2 = _mix(fill, style.accent if style.is_dark else "FFFFFF", 0.5)
        _grad(sp, fill, c2, 110)
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
        if getattr(style, "surface_alpha", 100) < 100:
            _apply_alpha(sp, style.surface_alpha)
    border = _panel_border(style)
    if border:
        sp.line.color.rgb = _rgb(border[0]); sp.line.width = border[1]
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    ss = getattr(style, "shadow_style", "soft")
    if ss == "glow":
        _glow(sp, style.accent)
    elif ss == "hard":
        _shadow(sp, blur=0, dist=48000, direction=2700000, alpha=42000,
                color=_mix(style.accent, "000000", 0.2))
    elif ss == "soft" and not style.is_dark:
        _shadow(sp)
    return sp


def _texture(slide, style: DeckStyle):
    """背景纹理层(在 _bg 之后、内容之前铺)· 撑不同艺术方向的"材质感"。"""
    tx = getattr(style, "texture", "none")
    if tx == "none":
        return
    W, H = int(SLIDE_W), int(SLIDE_H)
    if tx == "scanline":                       # 故障:横向扫描线
        col = _mix(style.bg, style.ink_body, 0.10 if style.is_dark else 0.06)
        step = int(Inches(0.18)); y = 0
        while y < H:
            _rect(slide, 0, y, W, Pt(1.1), col); y += step
    elif tx == "grid":                         # 手绘/蓝图:方格纸
        col = _mix(style.bg, style.ink_muted, 0.16)
        step = int(Inches(0.82)); x = step
        while x < W:
            _rect(slide, x, 0, Pt(0.75), H, col); x += step
        y = step
        while y < H:
            _rect(slide, 0, y, W, Pt(0.75), col); y += step
    elif tx == "dots":                         # 点阵
        col = _mix(style.bg, style.ink_muted, 0.22)
        step = int(Inches(0.95)); d = int(Inches(0.04)); y = step
        while y < H:
            x = step
            while x < W:
                _oval(slide, x, y, d, d, col); x += step
            y += step
    elif tx == "beams":                        # 玻璃:斜向柔光斑
        c1 = _mix(style.bg, style.accent, 0.55 if style.is_dark else 0.28)
        o1 = _oval(slide, int(Inches(-1.8)), int(Inches(-2.2)), int(Inches(7.4)), int(Inches(7.4)), c1)
        _apply_alpha(o1, 30)
        c2 = _mix(style.bg, getattr(style, "accent2", style.accent), 0.5)
        o2 = _oval(slide, int(Inches(9.2)), int(Inches(3.4)), int(Inches(6.2)), int(Inches(6.2)), c2)
        _apply_alpha(o2, 24)


# 当前生效字体(读系统字体库挑出来的 · 由 render 在开场设置)· None=退回 style 内置字体
_ACTIVE_FONTS: Optional[dict] = None


def set_active_fonts(fonts: Optional[dict]) -> None:
    global _ACTIVE_FONTS
    _ACTIVE_FONTS = dict(fonts) if fonts else None


def _faces(style: DeckStyle, display: bool):
    """(拉丁, 中日韩) 字体族 · display=大标题用更有设计感的那款。"""
    if _ACTIVE_FONTS:
        return _ACTIVE_FONTS["latin"], (_ACTIVE_FONTS["cjk_display"] if display else _ACTIVE_FONTS["cjk_body"])
    return style.font_en, style.font_cjk


def _set_font(run, style: DeckStyle, size_pt, color: str, bold: bool, display: bool = False):
    latin, cjk = _faces(style, display)
    f = run.font
    f.size = Pt(size_pt); f.bold = bold; f.name = latin
    f.color.rgb = _rgb(color)
    rPr = run._r.get_or_add_rPr()
    for tag, face in (("a:latin", latin), ("a:ea", cjk), ("a:cs", latin)):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", face)


def _tb(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return box, tf


def _para(tf, text, style, size, color, *, bold=False, align=PP_ALIGN.LEFT, first=False,
          space_after=8, line_spacing=1.14, level=0, marker=None, marker_color=None, tracking=False,
          display=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.line_spacing = line_spacing; p.space_after = Pt(space_after); p.level = level
    if marker:
        mr = p.add_run(); mr.text = f"{marker}   "
        _set_font(mr, style, size, marker_color or color, True)
    if tracking:                                   # 眉标字距(letter-spacing 近似)· 不解析 markdown
        tr = p.add_run(); tr.text = "  ".join(list(text or ""))
        _set_font(tr, style, size, color, bold, display=display)
        return p
    for seg, seg_bold in _inline_segments(text or ""):   # 正文:真解析 **粗体** 等 inline markdown
        if not seg:
            continue
        r = p.add_run(); r.text = seg
        _set_font(r, style, size, color, bold or seg_bold, display=display)
    return p


def _kicker_row(slide, style: DeckStyle, text: str, y, ink: Optional[str] = None):
    if not text:
        return
    ink = ink or style.accent
    txt = text.upper() if style.uppercase_kicker else text
    _rect(slide, MARGIN, y + Inches(0.07), Inches(0.30), Pt(4), ink)
    _, tf = _tb(slide, MARGIN + Inches(0.44), y, CONTENT_W - Inches(0.44), Inches(0.34))
    _para(tf, txt, style, style.pt_kicker, ink, bold=True, first=True, space_after=0, tracking=True)


def _header(slide, style: DeckStyle, kicker: Optional[str], title: Optional[str],
            subtitle: Optional[str] = None) -> Emu:
    """眉标 → hairline → 标题 → 副标题 · hairline 在标题之上(根治压字)· 返回统一正文起点。"""
    if kicker:
        _kicker_row(slide, style, kicker, Inches(0.58))
    hair_y = Inches(1.02)
    _rect(slide, MARGIN, hair_y, CONTENT_W, Pt(1), style.rule)
    _rect(slide, MARGIN, hair_y - Pt(1), Inches(1.15), Pt(3.2), style.accent)
    # 标题过长时收字号(防两行标题压到副标题)
    tsize = _fit_pt(title or "", style.pt_title, CONTENT_W, max_lines=1, min_pt=max(style.pt_heading, 18))
    _, tf = _tb(slide, MARGIN, Inches(1.22), CONTENT_W, Inches(0.95))
    _para(tf, title or "", style, tsize, style.ink_title, bold=True, first=True,
          space_after=2, line_spacing=1.02, display=True)
    if subtitle:
        _, tf2 = _tb(slide, MARGIN, Inches(2.02), CONTENT_W, Inches(0.44))
        _para(tf2, subtitle, style, style.pt_heading, style.ink_muted, first=True, space_after=0)
        return Inches(2.62)          # 有副标题 → 正文下移,别压副标题(卷七十九事故)
    return BODY_TOP


def _footer(slide, style: DeckStyle, page: Optional[int], footer_text: Optional[str]):
    ft = footer_text if footer_text is not None else style.footer
    _rect(slide, MARGIN, SLIDE_H - Inches(0.62), CONTENT_W, Pt(0.75), style.rule)
    _, tf = _tb(slide, MARGIN, SLIDE_H - Inches(0.52), CONTENT_W, Inches(0.3), MSO_ANCHOR.MIDDLE)
    _para(tf, ft or "", style, style.pt_footnote, style.ink_muted, first=True, space_after=0)
    if page is not None:
        _, tf2 = _tb(slide, SLIDE_W - MARGIN - Inches(1.0), SLIDE_H - Inches(0.52),
                     Inches(1.0), Inches(0.3), MSO_ANCHOR.MIDDLE)
        _para(tf2, f"{page:02d}", style, style.pt_footnote, style.accent,
              align=PP_ALIGN.RIGHT, bold=True, first=True, space_after=0)


def _watermark(slide, style: DeckStyle, text: str, color: str):
    """章节巨型数字/符号水印 · 右侧竖向锚定 · 低对比。"""
    _, tf = _tb(slide, SLIDE_W - Inches(6.2), Inches(1.1), Inches(5.7), Inches(5.6), MSO_ANCHOR.MIDDLE)
    _para(tf, text, style, 240, color, bold=True, first=True, align=PP_ALIGN.RIGHT,
          space_after=0, line_spacing=0.9, display=True)


def _pic_glyph(slide, cx, cy, sz, ink):
    """简单的"图片"象形:圆角相框 + 太阳 + 山峰(占位卡用 · 不引 icons 免循环依赖)。"""
    fr = _round(slide, cx - sz / 2, cy - sz / 2, sz, sz, ink, radius=0.12)
    fr.fill.background(); fr.line.color.rgb = _rgb(ink); fr.line.width = Pt(2.4); fr.shadow.inherit = False
    _oval(slide, cx - sz * 0.22, cy - sz * 0.26, sz * 0.16, sz * 0.16, ink)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, int(cx - sz * 0.34),
                                 int(cy + sz * 0.02), int(sz * 0.68), int(sz * 0.28))
    tri.fill.solid(); tri.fill.fore_color.rgb = _rgb(ink); tri.line.fill.background(); tri.shadow.inherit = False


def _fill_photo(slide, path, l, t, w, h):
    """把图裁剪填满 (l,t,w,h) 区域 · 居中裁 · 返回 picture shape。
    用于满版封面/hero(裁一点点、边到边更高级·且上面有 scrim+标题兜可读性)。"""
    pic = slide.shapes.add_picture(str(path), int(l), int(t))
    nw, nh = pic.width, pic.height
    box_ar = w / h
    img_ar = (nw / nh) if nh else box_ar
    if img_ar > box_ar:
        crop = (1 - box_ar / img_ar) / 2
        pic.crop_left = pic.crop_right = crop
    else:
        crop = (1 - img_ar / box_ar) / 2
        pic.crop_top = pic.crop_bottom = crop
    pic.left, pic.top, pic.width, pic.height = int(l), int(t), int(w), int(h)
    return pic


def _place_photo(slide, path, l, t, w, h, max_crop: float = 0.14):
    """智能放图(内容图页用):裁切量小(≤max_crop·两边共 ≤2×)→ 填满·边到边更高级;
    裁切量大(比例差很多·如竖图/方图塞进宽框)→ 等比装下(contain)居中·**不砍图上内容**。
    返回图片实际矩形 (l,t,w,h)(给图注定位)。 卷七十九续十七 · 治"竖图被裁掉标签"。"""
    pic = slide.shapes.add_picture(str(path), int(l), int(t))
    nw, nh = pic.width, pic.height
    box_ar = w / h
    img_ar = (nw / nh) if nh else box_ar
    if img_ar > box_ar:
        crop_frac = (1 - box_ar / img_ar) / 2      # 填满需裁左右
    else:
        crop_frac = (1 - img_ar / box_ar) / 2      # 填满需裁上下
    if crop_frac <= max_crop:
        if img_ar > box_ar:
            pic.crop_left = pic.crop_right = crop_frac
        else:
            pic.crop_top = pic.crop_bottom = crop_frac
        pic.left, pic.top, pic.width, pic.height = int(l), int(t), int(w), int(h)
        return int(l), int(t), int(w), int(h)
    # contain:等比装下·居中·一点不砍
    if img_ar > box_ar:
        dw = int(w); dh = int(round(w / img_ar))
    else:
        dh = int(h); dw = int(round(h * img_ar))
    dl = int(l) + (int(w) - dw) // 2
    dt = int(t) + (int(h) - dh) // 2
    pic.left, pic.top, pic.width, pic.height = dl, dt, dw, dh
    return dl, dt, dw, dh


def _image_cover(slide, path: Optional[Path], l, t, w, h, style: DeckStyle,
                 caption: Optional[str] = None, prompt: Optional[str] = None):
    if not path or not Path(path).exists():
        _round(slide, l, t, w, h, style.placeholder_fill, line=style.rule)
        cx = int(l) + int(w) / 2
        gy = int(t) + int(h) * 0.30
        _pic_glyph(slide, cx, gy, min(int(h) * 0.26, int(Inches(1.3))), style.placeholder_ink)
        _, tk = _tb(slide, int(l) + int(Inches(0.3)), int(t) + int(h) * 0.44,
                    int(w) - int(Inches(0.6)), int(Inches(0.34)))
        _para(tk, ("IMAGE PROMPT" if style.uppercase_kicker else "配图提示词"), style,
              style.pt_kicker, style.accent, bold=True, first=True, align=PP_ALIGN.CENTER,
              space_after=0, tracking=True)
        body = prompt or caption or "在此填写要生成的画面描述 · 或先取图再引用"
        _, tf = _tb(slide, int(l) + int(Inches(0.55)), int(t) + int(h) * 0.55,
                    int(w) - int(Inches(1.1)), int(h) * 0.4, MSO_ANCHOR.TOP)
        _para(tf, body, style, style.pt_body, style.placeholder_ink, first=True,
              align=PP_ALIGN.CENTER, space_after=0, line_spacing=1.25)
        return
    il, it, iw, ih = _place_photo(slide, path, l, t, w, h)
    if caption:
        # 图注跟着"图片实际底部"居中(contain 时图没占满框·别飘在半空/越界)
        cy = min(int(it) + int(ih) + int(Inches(0.08)), int(SLIDE_H) - int(Inches(0.62)))
        _, tf = _tb(slide, il, cy, iw, Inches(0.32))
        _para(tf, caption, style, style.pt_footnote, style.ink_muted, first=True,
              align=PP_ALIGN.CENTER, space_after=0)


def _bullets(slide, style: DeckStyle, bullets, l, t, w, h):
    """彩色方块 marker + 按【可用高度】auto-fit 收字号与段距 + 大行距。
    卷七十九续十九:治要点多/长时正文溢出 PPT 底边(BRO Fig1 伏隔核 8 条挤出页面)。
    做法:从基准字号往下试·估算每条换行行数 × 行高 + 段后间距·总高塞进 h 才停。"""
    n = len(bullets)
    if n == 0:
        return
    avail_in = float(h) / _EMU_PER_IN
    w_in = float(w) / _EMU_PER_IN
    LH = 1.16                                  # 行距倍数
    base = int(getattr(style, "pt_body", 18))
    sa_ratio, sa_sub = 0.60, 0.28              # 段后间距 / 字号
    size = max(base, 12)
    for cand in range(base, 11, -1):
        size = cand
        cap_u = max((w_in - 0.30) * 72.0 / cand, 1.0)   # 每行能塞的 units(扣掉 marker 缩进)
        total = 0.0
        for b in bullets:
            sub = b.get("level", 0) > 0
            es = cand if not sub else max(cand - 2, 11)
            u = _text_units(b.get("text", ""))
            lines = max(1, int(u / cap_u + 0.999))
            total += lines * es / 72.0 * LH + es / 72.0 * (sa_ratio if not sub else sa_sub)
        if total <= avail_in:
            break
    _, tf = _tb(slide, l, t, w, h)
    for i, b in enumerate(bullets):
        lvl = b.get("level", 0)
        sub = lvl > 0
        es = size if not sub else max(size - 2, 11)
        sa = max(3, round(es * (sa_ratio if not sub else sa_sub)))
        _para(
            tf, b["text"], style, es,
            style.ink_body if not sub else style.ink_muted,
            first=(i == 0), space_after=sa,
            line_spacing=LH, level=lvl,
            marker=("▪" if not sub else "·"),
            marker_color=(style.accent if not sub else style.ink_muted),
        )
